
from flask import Flask, jsonify
from google.oauth2 import service_account
from googleapiclient.discovery import build
import io
import fitz
import docx
import openpyxl
from pptx import Presentation
import os

app = Flask(__name__)

SERVICE_ACCOUNT_FILE = 'service_account.json'
SCOPES = ['https://www.googleapis.com/auth/drive']
creds = service_account.Credentials.from_service_account_file(SERVICE_ACCOUNT_FILE, scopes=SCOPES)
drive_service = build('drive', 'v3', credentials=creds)

FOLDER_ID = "1jHDjMkJeRbuMZiJW8zM-pVjhx0jDUfPY"

def list_all_files(folder_id):
    files = []
    queue = [folder_id]
    while queue:
        current = queue.pop(0)
        results = drive_service.files().list(
            q=f"'{current}' in parents and trashed=false",
            fields="files(id, name, mimeType, webContentLink, webViewLink)",
            supportsAllDrives=True,
            includeItemsFromAllDrives=True
        ).execute()
        for file in results.get('files', []):
            if file['mimeType'] == 'application/vnd.google-apps.folder':
                queue.append(file['id'])
            else:
                files.append(file)
    return files

def extract_text(file_id, mime_type):
    try:
        if mime_type == 'application/pdf':
            media = drive_service.files().get_media(fileId=file_id)
            response, content = drive_service._http.request(media.uri)
            doc = fitz.open(stream=content, filetype="pdf")
            return "".join(page.get_text() for page in doc)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            media = drive_service.files().get_media(fileId=file_id)
            response, content = drive_service._http.request(media.uri)
            with io.BytesIO(content) as f:
                doc = docx.Document(f)
                return "\n".join([para.text for para in doc.paragraphs])
        elif mime_type == 'application/vnd.google-apps.document':
            media = drive_service.files().export_media(
                fileId=file_id,
                mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
            response, content = drive_service._http.request(media.uri)
            with io.BytesIO(content) as f:
                doc = docx.Document(f)
                return "\n".join([para.text for para in doc.paragraphs])
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            media = drive_service.files().get_media(fileId=file_id)
            response, content = drive_service._http.request(media.uri)
            with io.BytesIO(content) as f:
                wb = openpyxl.load_workbook(f, data_only=True)
                return "\n".join(" ".join(str(cell) if cell else "" for cell in row) for sheet in wb.worksheets for row in sheet.iter_rows(values_only=True))
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            media = drive_service.files().get_media(fileId=file_id)
            response, content = drive_service._http.request(media.uri)
            with io.BytesIO(content) as f:
                prs = Presentation(f)
                return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    except:
        return ""

@app.route('/fetch-drive-documents', methods=['GET'])
def fetch_drive_documents():
    files = list_all_files(FOLDER_ID)
    supported_types = {
        'application/pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'application/vnd.google-apps.document',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    }

    documents = []
    for file in files:
        if file['mimeType'] not in supported_types:
            continue

        text = extract_text(file['id'], file['mimeType'])
        if text and text.strip():
            documents.append({
                "file_id": file['id'],
                "file_name": file['name'],
                "file_type": file['mimeType'],
                "text": text,
                "download_url": file.get('webContentLink') or f"https://drive.google.com/uc?id={file['id']}",
                "view_url": file.get('webViewLink')
            })
    return jsonify({"documents": documents})

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port)
